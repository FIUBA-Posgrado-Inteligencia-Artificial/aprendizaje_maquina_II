"""
Extract slides from a PPTX or PDF file into Markdown + images.

Usage:
    uv run python extract_presentation.py <input.pptx|input.pdf> <output_dir>

Output:
    <output_dir>/slides.md      — Markdown with one section per slide
    <output_dir>/images/        — Extracted images referenced in the Markdown
"""

import os
import re
import sys
from pathlib import Path


# ── PPTX extraction ──────────────────────────────────────────────────────────

def _run_fragment(run):
    """Return (text, is_bold) for a pptx run."""
    bold = run.font.bold
    return run.text, bool(bold)


def _para_to_md(para):
    """Convert a pptx paragraph to a Markdown string (inline formatting only)."""
    parts = []
    for run in para.runs:
        text, bold = _run_fragment(run)
        if not text.strip():
            continue
        if bold:
            text = f"**{text.strip()}**"
        parts.append(text)
    return "".join(parts)


_SKIP_PLACEHOLDER_IDX = {11, 12, 13, 14}  # footer, slide number, date, etc.

# Content placeholders (idx=1) have bullets defined in the master/layout,
# so buChar/buAutoNum flags are absent — use para.level directly.
_BODY_PLACEHOLDER_IDX = {1, 2, 3, 4, 5, 6, 7, 8, 9, 10}


def _shape_text_to_md(shape):
    """Return list of Markdown lines for a text shape."""
    is_body = False
    try:
        ph = shape.placeholder_format
    except ValueError:
        ph = None
    if ph is not None:
        if ph.idx in _SKIP_PLACEHOLDER_IDX:
            return []
        if ph.idx in _BODY_PLACEHOLDER_IDX:
            is_body = True

    lines = []
    for para in shape.text_frame.paragraphs:
        text = _para_to_md(para).strip()
        if not text:
            continue
        level = para.level if para.level is not None else 0
        if is_body:
            indent = "  " * level
            lines.append(f"{indent}* {text}")
        else:
            lines.append(text)
    return lines


def extract_pptx(input_path: Path, output_dir: Path) -> str:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    images_dir = output_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(input_path)
    slide_sections = []
    img_counter = 1

    for slide_num, slide in enumerate(prs.slides, start=1):
        title_text = ""
        body_lines = []
        slide_images = []

        # Separate title placeholder from body shapes
        title_shape = None
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_text = title_shape.text.strip()

        for shape in slide.shapes:
            if shape == title_shape:
                continue

            # Images / pictures
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = shape.image
                ext = img.ext
                img_name = f"img-{img_counter:03d}.{ext}"
                img_path = images_dir / img_name
                img_path.write_bytes(img.blob)
                slide_images.append(f"images/{img_name}")
                img_counter += 1
                continue

            # Text boxes / placeholders
            if shape.has_text_frame:
                lines = _shape_text_to_md(shape)
                body_lines.extend(lines)

        # Build the slide markdown
        parts = []
        if title_text:
            parts.append(f"## Diapositiva {slide_num}: {title_text}")
        else:
            parts.append(f"## Diapositiva {slide_num}")

        parts.extend(body_lines)

        for img_ref in slide_images:
            img_name = Path(img_ref).stem
            parts.append(f"![{img_name}]({img_ref})")

        slide_sections.append("\n\n".join(parts) if len(parts) > 1 else parts[0])

    md = "\n\n---\n\n".join(slide_sections)
    return md


# ── PDF extraction ────────────────────────────────────────────────────────────

def extract_pdf(input_path: Path, output_dir: Path) -> str:
    import fitz  # pymupdf

    images_dir = output_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    doc = fitz.open(input_path)
    slide_sections = []
    img_counter = 1

    for page_num, page in enumerate(doc, start=1):
        lines = []
        slide_images = []

        # Extract text blocks with basic structure
        blocks = page.get_text("dict")["blocks"]
        for block in blocks:
            if block["type"] == 0:  # text block
                for line in block["lines"]:
                    spans = line["spans"]
                    if not spans:
                        continue
                    text_parts = []
                    for span in spans:
                        text = span["text"].strip()
                        if not text:
                            continue
                        flags = span["flags"]
                        is_bold = bool(flags & 16)
                        if is_bold:
                            text = f"**{text}**"
                        text_parts.append(text)
                    combined = " ".join(text_parts).strip()
                    # Skip lines that are just slide numbers
                    if combined and not re.fullmatch(r"\d+", combined):
                        # Convert PDF bullet chars to markdown
                        combined = re.sub(r"^[•◦▪▸►‣–\-]\s*", "* ", combined)
                        lines.append(combined)

            elif block["type"] == 1:  # image block
                # Extract the image
                img_index = block.get("number", img_counter)
                img_list = page.get_images(full=True)
                # Match by position heuristic — just save page images in order
                pass

        # Extract images from the page
        for img_info in page.get_images(full=True):
            xref = img_info[0]
            base_image = doc.extract_image(xref)
            img_bytes = base_image["image"]
            ext = base_image["ext"]
            img_name = f"img-{img_counter:03d}.{ext}"
            img_path = images_dir / img_name
            img_path.write_bytes(img_bytes)
            slide_images.append((f"images/{img_name}", img_name))
            img_counter += 1

        # First non-empty line is the slide "title"
        title = ""
        body = []
        if lines:
            title = lines[0]
            body = lines[1:]

        parts = []
        if title:
            # Remove bold markers from title for the heading
            clean_title = re.sub(r"\*\*(.+?)\*\*", r"\1", title)
            parts.append(f"## Diapositiva {page_num}: {clean_title}")
        else:
            parts.append(f"## Diapositiva {page_num}")

        parts.extend(body)

        for img_ref, img_name in slide_images:
            parts.append(f"![{img_name}]({img_ref})")

        slide_sections.append("\n\n".join(parts) if len(parts) > 1 else parts[0])

    md = "\n\n---\n\n".join(slide_sections)
    return md


# ── Entry point ───────────────────────────────────────────────────────────────

def main():
    if len(sys.argv) < 3:
        print("Uso: uv run python extract_presentation.py <archivo.pptx|archivo.pdf> <directorio_salida>")
        sys.exit(1)

    input_path = Path(sys.argv[1]).resolve()
    output_dir = Path(sys.argv[2]).resolve()

    if not input_path.exists():
        print(f"Error: no se encontró el archivo {input_path}")
        sys.exit(1)

    output_dir.mkdir(parents=True, exist_ok=True)

    ext = input_path.suffix.lower()
    if ext == ".pptx":
        md = extract_pptx(input_path, output_dir)
    elif ext == ".pdf":
        md = extract_pdf(input_path, output_dir)
    else:
        print(f"Error: formato no soportado '{ext}'. Usá .pptx o .pdf")
        sys.exit(1)

    out_md = output_dir / "slides.md"
    out_md.write_text(md, encoding="utf-8")
    print(f"Markdown generado en: {out_md}")
    print(f"Imágenes guardadas en: {output_dir / 'images'}/")


if __name__ == "__main__":
    main()
